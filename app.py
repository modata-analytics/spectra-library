# Author: modata-analytics

from pathlib import Path
from io import BytesIO
import re
import tempfile

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
from pyteomics import mzxml
from pptx import Presentation
from pptx.util import Inches


st.set_page_config(
    page_title="Microbial Reference Library",
    page_icon="🧪",
    layout="wide"
)

DATA_DIR = Path("data")
METADATA_FILE = DATA_DIR / "metadata.csv"
SPECTRA_DIR = DATA_DIR / "spectra"

DISPLAY_COLUMN_NAMES = {
    "microbe": "Microbe",
    "organism_type": "Organism Type",
    "replicate_id": "Replicate ID",
    "replicate_type": "Replicate Type",
    "strain": "Strain",
    "condition": "Condition",
    "gram_status": "Gram Status",
    "analyte_category": "Analyte Category",
    "reference_csv": "Reference CSV",
    "reference_image": "Reference Image",
    "reference_mzxml": "Reference mzXML",
    "notes": "Notes",
    "keyword_tags": "Keyword Tags",
    "n_records": "Record Count",
    "mz": "m/z",
    "intensity": "Intensity",
}


def normalize_search_text(text) -> str:
    text = str(text).lower().strip()
    text = text.replace("-", " ")
    text = text.replace("_", " ")
    text = re.sub(r"\s+", " ", text)
    return text


@st.cache_data
def load_metadata() -> pd.DataFrame:
    if not METADATA_FILE.exists():
        return pd.DataFrame()

    df = pd.read_csv(METADATA_FILE)
    df = df.fillna("")

    if "organism_type" not in df.columns:
        df["organism_type"] = "Bacterium"

    if "keyword_tags" not in df.columns:
        df["keyword_tags"] = ""

    if "reference_mzxml" not in df.columns:
        df["reference_mzxml"] = ""

    if "organism_type" in df.columns:
        df["organism_type"] = (
            df["organism_type"]
            .astype(str)
            .str.strip()
            .replace({
                "bacterium": "Bacterium",
                "bacteria": "Bacterium",
                "Bacteria": "Bacterium",
                "Bacterium": "Bacterium",
                "yeast": "Yeast",
                "Yeast": "Yeast",
                "fungus": "Fungal",
                "fungal": "Fungal",
                "Fungus": "Fungal",
                "Fungal": "Fungal",
            })
        )

    if "analyte_category" in df.columns:
        df["analyte_category"] = (
            df["analyte_category"]
            .astype(str)
            .str.strip()
            .replace({
                "lipid_A": "Lipid A",
                "lipid a": "Lipid A",
                "Lipid A": "Lipid A",
                "lipid": "Lipid A",
                "Lipid": "Lipid A",
                "los": "LOS Like",
                "LOS": "LOS Like",
                "LOS_like": "LOS Like",
                "los_like": "LOS Like",
                "los like": "LOS Like",
                "LOS Like": "LOS Like",
                "not_applicable": "Cardiolipin Profile",
                "cardio": "Cardiolipin Profile",
                "Cardio": "Cardiolipin Profile",
                "cardiolipin_profile": "Cardiolipin Profile",
                "cardiolipin profile": "Cardiolipin Profile",
                "Cardiolipin Profile": "Cardiolipin Profile",
                "sterol profile": "Sterol Profile",
                "Sterol Profile": "Sterol Profile",
                "sterol_profile": "Sterol Profile",
            })
        )

    if "gram_status" in df.columns:
        df["gram_status"] = (
            df["gram_status"]
            .astype(str)
            .str.strip()
            .replace({
                "gram-negative": "Gram-negative",
                "gram negative": "Gram-negative",
                "Gram-negative": "Gram-negative",
                "neg": "Gram-negative",
                "Neg": "Gram-negative",
                "gram-positive": "Gram-positive",
                "gram positive": "Gram-positive",
                "Gram-positive": "Gram-positive",
                "pos": "Gram-positive",
                "Pos": "Gram-positive",
                "gram-variable": "Gram-variable",
                "gram variable": "Gram-variable",
                "Gram-variable": "Gram-variable",
                "na": "Not Applicable",
                "N/A": "Not Applicable",
                "n/a": "Not Applicable",
                "not applicable": "Not Applicable",
                "Not Applicable": "Not Applicable",
            })
        )

    if "replicate_type" in df.columns:
        df["replicate_type"] = (
            df["replicate_type"]
            .astype(str)
            .str.strip()
            .replace({
                "technical": "Technical",
                "Technical": "Technical",
                "tech": "Technical",
                "Tech": "Technical",
                "biological": "Biological",
                "Biological": "Biological",
                "bio": "Biological",
                "Bio": "Biological",
            })
        )

    return df


def prettify_dataframe_columns(df: pd.DataFrame) -> pd.DataFrame:
    renamed_columns = {
        col: DISPLAY_COLUMN_NAMES.get(col, col.replace("_", " ").title())
        for col in df.columns
    }
    return df.rename(columns=renamed_columns)


def build_search_mask(df: pd.DataFrame, search_term: str) -> pd.Series:
    search_cols = [
        "microbe",
        "organism_type",
        "strain",
        "notes",
        "condition",
        "analyte_category",
        "gram_status",
        "keyword_tags",
    ]
    valid_cols = [col for col in search_cols if col in df.columns]

    normalized_term = normalize_search_text(search_term)
    search_tokens = normalized_term.split()

    def row_matches(row):
        combined_text = " ".join(
            normalize_search_text(row[col]) for col in valid_cols
        )
        return all(token in combined_text for token in search_tokens)

    return df.apply(row_matches, axis=1)


def read_mzxml_scans(file_obj_or_path):
    scans = []
    with mzxml.MzXML(file_obj_or_path) as reader:
        for scan in reader:
            mzs = np.asarray(scan.get("m/z array", []), dtype=float)
            intensities = np.asarray(scan.get("intensity array", []), dtype=float)
            ms_level = scan.get("msLevel", None)

            if mzs.size > 0 and intensities.size > 0:
                scans.append({
                    "num": scan.get("num"),
                    "ms_level": ms_level,
                    "mzs": mzs,
                    "intensities": intensities,
                    "tic": float(np.sum(intensities)),
                })
    return scans


def choose_representative_scan(scans, ms_level=1):
    valid = [s for s in scans if s["ms_level"] == ms_level]
    if not valid:
        valid = scans
    if not valid:
        return None
    return max(valid, key=lambda s: s["tic"])


def normalize_intensity(intensities):
    if intensities.size == 0:
        return intensities
    max_i = np.max(intensities)
    if max_i == 0:
        return intensities
    return intensities / max_i


def pick_peaks(mzs, intensities, rel_height=0.10, min_distance_mz=0.5):
    if len(mzs) < 3:
        return mzs, intensities

    picked_mz = []
    picked_i = []
    threshold = rel_height * np.max(intensities)

    for i in range(1, len(mzs) - 1):
        if (
            intensities[i] >= threshold
            and intensities[i] >= intensities[i - 1]
            and intensities[i] >= intensities[i + 1]
        ):
            if not picked_mz or (mzs[i] - picked_mz[-1]) >= min_distance_mz:
                picked_mz.append(mzs[i])
                picked_i.append(intensities[i])
            else:
                if intensities[i] > picked_i[-1]:
                    picked_mz[-1] = mzs[i]
                    picked_i[-1] = intensities[i]

    return np.array(picked_mz), np.array(picked_i)


def annotate_top_peaks(ax, peak_mz, peak_i, n=10):
    if len(peak_mz) == 0:
        return

    idx = np.argsort(peak_i)[-n:]
    idx = idx[np.argsort(peak_mz[idx])]

    x_min, x_max = ax.get_xlim()
    y_min, y_max = ax.get_ylim()

    y_offsets = [0.04, 0.08, 0.12, 0.16]
    x_offsets = [-12, -6, 6, 12]

    for j, i in enumerate(idx):
        x = float(peak_mz[i])
        y = float(peak_i[i])

        x = min(max(x, x_min + 10), x_max - 10)

        label_x = x + x_offsets[j % len(x_offsets)]
        label_y = min(y + y_offsets[j % len(y_offsets)], y_max * 0.96)

        ax.annotate(
            f"{peak_mz[i]:.1f}",
            xy=(x, y),
            xytext=(label_x, label_y),
            textcoords="data",
            fontsize=8,
            rotation=35,
            ha="center",
            va="bottom",
            clip_on=True,
            arrowprops=dict(arrowstyle="-", lw=0.5, shrinkA=0, shrinkB=0),
        )


def plot_mzxml_spectrum(mzs, intensities, peak_mz, peak_i, title):
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.plot(mzs, intensities, linewidth=1)

    if len(peak_mz) > 0:
        ax.scatter(peak_mz, peak_i, s=18)

    ax.set_title(title)
    ax.set_xlabel("m/z")
    ax.set_ylabel("Normalized Intensity")
    ax.set_ylim(0, max(1.05, np.max(intensities) * 1.08))

    x_min = np.min(mzs)
    x_max = np.max(mzs)
    x_pad = (x_max - x_min) * 0.05
    ax.set_xlim(x_min - x_pad, x_max + x_pad)

    if len(peak_mz) > 0:
        annotate_top_peaks(ax, peak_mz, peak_i, n=10)

    fig.tight_layout()
    return fig


def spectrum_from_mzxml(file_obj_or_path, rel_height=0.10, min_distance_mz=0.5):
    scans = read_mzxml_scans(file_obj_or_path)
    rep = choose_representative_scan(scans, ms_level=1)
    if rep is None:
        return None

    mzs = rep["mzs"]
    intensities = normalize_intensity(rep["intensities"])
    peak_mz, peak_i = pick_peaks(
        mzs,
        intensities,
        rel_height=rel_height,
        min_distance_mz=min_distance_mz,
    )

    return {
        "scan_num": rep["num"],
        "mzs": mzs,
        "intensities": intensities,
        "peak_mz": peak_mz,
        "peak_i": peak_i,
    }


@st.cache_data
def load_reference_library(metadata_df: pd.DataFrame):
    refs = []

    for _, row in metadata_df.iterrows():
        mzxml_path = str(row.get("reference_mzxml", "")).strip()
        if not mzxml_path:
            continue

        path = Path(mzxml_path)
        if not path.exists():
            continue

        parsed = spectrum_from_mzxml(str(path))
        if parsed is not None:
            refs.append({
                "microbe": row.get("microbe", path.stem),
                "organism_type": row.get("organism_type", ""),
                "gram_status": row.get("gram_status", ""),
                "analyte_category": row.get("analyte_category", ""),
                "replicate_id": row.get("replicate_id", ""),
                "file": str(path),
                **parsed,
            })

    return refs


def get_top_peaks(peak_mz, peak_i, top_n=25):
    if len(peak_mz) == 0:
        return np.array([]), np.array([])

    idx = np.argsort(peak_i)[-top_n:]
    idx = idx[np.argsort(peak_mz[idx])]
    return peak_mz[idx], peak_i[idx]


def infer_profile_class(query_spec):
    peak_mz = query_spec["peak_mz"]
    peak_i = query_spec["peak_i"]

    if len(peak_mz) == 0:
        return "Unknown"

    top_mz, top_i = get_top_peaks(peak_mz, peak_i, top_n=12)

    strong_high_mass = np.sum((top_mz >= 1600) & (top_i >= 0.25))
    strong_mid_mass = np.sum((top_mz >= 1200) & (top_mz <= 1500) & (top_i >= 0.25))

    if strong_high_mass >= 2:
        return "Gram-negative"
    if strong_mid_mass >= 2 and strong_high_mass == 0:
        return "Gram-positive"
    return "Unknown"


def peak_match_similarity(query_mz, query_i, ref_mz, ref_i, tolerance=1.0):
    if len(query_mz) == 0 or len(ref_mz) == 0:
        return 0.0

    score = 0.0
    used_ref = set()

    for q_mz, q_i in zip(query_mz, query_i):
        diffs = np.abs(ref_mz - q_mz)
        best_idx = np.argmin(diffs)
        best_diff = diffs[best_idx]

        if best_diff <= tolerance and best_idx not in used_ref:
            closeness = 1.0 - (best_diff / tolerance)
            score += float(q_i) * float(ref_i[best_idx]) * closeness
            used_ref.add(best_idx)

    norm = np.linalg.norm(query_i) * np.linalg.norm(ref_i)
    if norm == 0:
        return 0.0

    return score / norm


def category_adjustment(query_class, ref_row):
    gram_status = str(ref_row.get("gram_status", ""))
    analyte_category = str(ref_row.get("analyte_category", ""))

    if query_class == "Gram-negative":
        if analyte_category in ["Lipid A", "LOS Like"] or gram_status == "Gram-negative":
            return 1.15
        if analyte_category == "Cardiolipin Profile" or gram_status == "Gram-positive":
            return 0.45
        if analyte_category == "Sterol Profile":
            return 0.35

    if query_class == "Gram-positive":
        if analyte_category == "Cardiolipin Profile" or gram_status == "Gram-positive":
            return 1.15
        if analyte_category in ["Lipid A", "LOS Like"] or gram_status == "Gram-negative":
            return 0.45
        if analyte_category == "Sterol Profile":
            return 0.35

    return 1.0


def compare_to_library(query_spec, refs, top_n=5, tolerance=1.0):
    if not refs:
        return pd.DataFrame()

    query_top_mz, query_top_i = get_top_peaks(query_spec["peak_mz"], query_spec["peak_i"], top_n=25)
    query_class = infer_profile_class(query_spec)

    rows = []
    for ref in refs:
        ref_top_mz, ref_top_i = get_top_peaks(ref["peak_mz"], ref["peak_i"], top_n=25)

        base_score = peak_match_similarity(
            query_top_mz,
            query_top_i,
            ref_top_mz,
            ref_top_i,
            tolerance=tolerance,
        )

        adjusted_score = base_score * category_adjustment(query_class, ref)

        rows.append({
            "Microbe": ref["microbe"],
            "Organism Type": ref["organism_type"],
            "Gram Status": ref["gram_status"],
            "Analyte Category": ref["analyte_category"],
            "Replicate ID": ref["replicate_id"],
            "Reference File": Path(ref["file"]).name,
            "Similarity": round(adjusted_score, 4),
            "Base Score": round(base_score, 4),
        })

    out = pd.DataFrame(rows).sort_values("Similarity", ascending=False).head(top_n)
    return out


def create_ppt_export(selected_row, parsed_spectrum):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = selected_row.get("microbe", "Spectrum Export")

    metadata_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(3.2), Inches(2.8))
    metadata_frame = metadata_box.text_frame
    metadata_frame.text = (
        f"Organism Type: {selected_row.get('organism_type', '')}\n"
        f"Gram Status: {selected_row.get('gram_status', '')}\n"
        f"Analyte Category: {selected_row.get('analyte_category', '')}\n"
        f"Strain: {selected_row.get('strain', '')}\n"
        f"Condition: {selected_row.get('condition', '')}\n"
        f"Replicate ID: {selected_row.get('replicate_id', '')}"
    )

    fig = plot_mzxml_spectrum(
        parsed_spectrum["mzs"],
        parsed_spectrum["intensities"],
        parsed_spectrum["peak_mz"],
        parsed_spectrum["peak_i"],
        selected_row.get("microbe", "Spectrum Export"),
    )

    tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(tmp_img.name, dpi=300, bbox_inches="tight")
    plt.close(fig)

    slide.shapes.add_picture(tmp_img.name, Inches(3.8), Inches(0.9), width=Inches(5.8))

    if len(parsed_spectrum["peak_mz"]) > 0:
        peak_df = pd.DataFrame({
            "mz": parsed_spectrum["peak_mz"],
            "intensity": parsed_spectrum["peak_i"],
        }).sort_values("intensity", ascending=False).head(10)

        peak_text = "Top Peaks:\n" + "\n".join(
            [f"{row.mz:.1f} | {row.intensity:.3f}" for row in peak_df.itertuples()]
        )

        peak_box = slide.shapes.add_textbox(Inches(0.4), Inches(4.0), Inches(3.0), Inches(2.2))
        peak_frame = peak_box.text_frame
        peak_frame.text = peak_text

    footer_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.8), Inches(6), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Developed by modata-analytics | https://github.com/modata-analytics"

    tmp_ppt = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_ppt.name)
    return tmp_ppt.name


def render_library_browser(metadata: pd.DataFrame):
    st.subheader("Library Browser")

    with st.sidebar:
        st.header("Library Filters")

        search_term = st.text_input(
            "Search Microbe / Strain / Notes / Keywords",
            value="",
            placeholder="Type keyword...",
            key="library_search",
            help="Search across organism names, notes, strain fields, and keyword tags with flexible partial matching.",
        )

        microbe_options = [""] + sorted(
            [x for x in metadata["microbe"].dropna().unique() if str(x).strip()]
        )
        selected_microbe = st.selectbox(
            "Select Microbe",
            options=microbe_options,
            index=0,
            format_func=lambda x: "All" if x == "" else x,
            key="library_microbe",
            help="Filter the library to one specific organism while leaving the rest of the filters unchanged.",
        )

        organism_options = [""] + sorted(
            [x for x in metadata["organism_type"].dropna().unique() if str(x).strip()]
        )
        selected_organism_type = st.selectbox(
            "Organism Type",
            options=organism_options,
            index=0,
            format_func=lambda x: "All" if x == "" else x,
            key="library_organism_type",
            help="Restrict the view to broad biological groups such as Bacterium, Yeast, or Fungal.",
        )

        gram_options = [""] + sorted(
            [x for x in metadata["gram_status"].dropna().unique() if str(x).strip()]
        )
        selected_gram = st.selectbox(
            "Gram Status",
            options=gram_options,
            index=0,
            format_func=lambda x: "All" if x == "" else x,
            key="library_gram_status",
            help="Narrow results by Gram classification or related status labels in the metadata.",
        )

        analyte_options = [""] + sorted(
            [x for x in metadata["analyte_category"].dropna().unique() if str(x).strip()]
        )
        selected_analyte = st.selectbox(
            "Analyte Category",
            options=analyte_options,
            index=0,
            format_func=lambda x: "All" if x == "" else x,
            key="library_analyte_category",
            help="Filter spectra by the profile type used for comparison, such as Lipid A or Cardiolipin Profile.",
        )

        replicate_options = [""] + sorted(
            [x for x in metadata["replicate_type"].dropna().unique() if str(x).strip()]
        )
        selected_replicate_type = st.selectbox(
            "Replicate Type",
            options=replicate_options,
            index=0,
            format_func=lambda x: "All" if x == "" else x,
            key="library_replicate_type",
            help="Focus the results on a replicate class such as Technical or Biological.",
        )

    filtered = metadata.copy()

    if search_term.strip():
        filtered = filtered[build_search_mask(filtered, search_term)]

    if selected_microbe != "":
        filtered = filtered[filtered["microbe"] == selected_microbe]

    if selected_organism_type != "":
        filtered = filtered[filtered["organism_type"] == selected_organism_type]

    if selected_gram != "":
        filtered = filtered[filtered["gram_status"] == selected_gram]

    if selected_analyte != "":
        filtered = filtered[filtered["analyte_category"] == selected_analyte]

    if selected_replicate_type != "":
        filtered = filtered[filtered["replicate_type"] == selected_replicate_type]

    if filtered.empty:
        st.error("No matching records were found. Check the search term or selected filters.")
        return

    filtered = filtered.reset_index(drop=True)
    filtered["display_label"] = filtered.apply(
        lambda row: f"{row['microbe']} | {row['replicate_id']} | {row['replicate_type']}",
        axis=1,
    )

    if selected_microbe:
        st.markdown(f"### {selected_microbe}")
    else:
        st.markdown("### Filtered Library Results")

    top_row = filtered.iloc[0]
    col1, col2 = st.columns([1, 2])

    with col1:
        st.markdown("### Metadata")
        st.write(f"Organism Type: {top_row['organism_type']}")
        st.write(f"Gram Status: {top_row['gram_status']}")
        st.write(f"Analyte Category: {top_row['analyte_category']}")
        st.write(f"Strain: {top_row['strain'] or 'Not provided'}")
        st.write(f"Condition: {top_row['condition'] or 'Not provided'}")
        st.write(f"Notes: {top_row['notes'] or 'None'}")
        st.write(f"Keyword Tags: {top_row['keyword_tags'] or 'None'}")

        if top_row["analyte_category"] == "Cardiolipin Profile":
            st.info(
                "This organism does not have a Lipid A profile, but it can still be "
                "distinguished within the library using its cardiolipin profile."
            )
        elif top_row["analyte_category"] == "Sterol Profile":
            st.info(
                "This organism is represented through a sterol-associated profile rather than "
                "a Lipid A profile."
            )

    with col2:
        st.markdown("### Matching Entries")
        display_cols = [
            "microbe",
            "organism_type",
            "replicate_id",
            "replicate_type",
            "strain",
            "condition",
            "gram_status",
            "analyte_category",
            "keyword_tags",
            "reference_mzxml",
            "reference_image",
        ]
        cols_present = [c for c in display_cols if c in filtered.columns]
        matching_entries_df = prettify_dataframe_columns(filtered[cols_present].copy())
        st.dataframe(matching_entries_df, use_container_width=True)

    st.markdown("---")
    st.markdown("### Reference Spectrum")

    selected_label = st.selectbox(
        "Select Entry",
        filtered["display_label"].tolist(),
        key="library_entry_select",
        help="Choose a specific entry to display its stored reference spectrum and export options.",
    )

    selected_row = filtered[filtered["display_label"] == selected_label].iloc[0]
    reference_mzxml = str(selected_row.get("reference_mzxml", "")).strip()

    left, right = st.columns([2, 1])

    with left:
        if reference_mzxml and Path(reference_mzxml).exists():
            parsed = spectrum_from_mzxml(reference_mzxml)

            if parsed is not None:
                fig = plot_mzxml_spectrum(
                    parsed["mzs"],
                    parsed["intensities"],
                    parsed["peak_mz"],
                    parsed["peak_i"],
                    f"{selected_row['microbe']} | {selected_row['replicate_id']} | {selected_row['replicate_type']}",
                )
                st.pyplot(fig)

                with st.expander("Show Peak Table"):
                    peak_df = pd.DataFrame({
                        "mz": parsed["peak_mz"],
                        "intensity": parsed["peak_i"],
                    }).sort_values("intensity", ascending=False)
                    peak_df = prettify_dataframe_columns(peak_df)
                    st.dataframe(peak_df, use_container_width=True)

                ppt_path = create_ppt_export(selected_row.to_dict(), parsed)
                with open(ppt_path, "rb") as f:
                    st.download_button(
                        label="Export This Spectrum to PowerPoint",
                        data=f,
                        file_name=f"{normalize_search_text(selected_row['microbe']).replace(' ', '_')}_spectrum.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="library_ppt_export",
                        help="Download a one-slide PowerPoint summary containing the selected spectrum and key metadata.",
                    )
            else:
                st.warning("No usable scans were found in the reference mzXML file.")
        else:
            st.warning("No reference mzXML file found for this entry.")

    with right:
        image_path = str(selected_row.get("reference_image", "")).strip()
        if image_path and Path(image_path).exists():
            st.image(image_path, caption=f"{selected_row['microbe']} spectrum image")
        else:
            st.info("No spectrum image found for this entry.")

    st.markdown("---")
    st.markdown("### Library Summary")

    summary = (
        metadata.groupby(["microbe", "organism_type", "gram_status", "analyte_category"])
        .size()
        .reset_index(name="n_records")
        .sort_values(["microbe", "analyte_category"])
    )
    summary = prettify_dataframe_columns(summary)
    st.dataframe(summary, use_container_width=True)


def render_mzxml_matcher(metadata: pd.DataFrame):
    st.subheader("mzXML Upload and Matching")
    st.markdown("Upload an mzXML file or choose one from the library to preview labeled peaks and identify closest matches.")

    input_col, settings_col = st.columns([2, 1])

    with input_col:
        mode = st.radio(
            "Choose Input Source",
            ["Upload mzXML", "Select Library mzXML"],
            horizontal=True,
            help="Choose whether to analyze a new uploaded file or inspect an mzXML file already stored in the reference library.",
        )

        uploaded = None
        selected_path = None

        if mode == "Upload mzXML":
            uploaded = st.file_uploader(
                "Drag and drop an mzXML file",
                type=["mzxml", "mzXML"],
                key="mzxml_upload",
                help="Upload a single mzXML spectrum file for visualization, peak picking, and library comparison.",
            )
        else:
            metadata_paths = [
                str(p).strip()
                for p in metadata.get("reference_mzxml", pd.Series(dtype=str)).tolist()
                if str(p).strip() and Path(str(p).strip()).exists()
            ]
            metadata_paths = sorted(list(set(metadata_paths)))
            chosen = st.selectbox(
                "Choose Reference mzXML File",
                options=metadata_paths if metadata_paths else [""],
                key="mzxml_library_select",
                help="Select an mzXML file already linked in the metadata library to preview it in the matcher.",
            )
            if chosen:
                selected_path = Path(chosen)

    with settings_col:
        rel_height = st.slider(
            "Peak Threshold",
            0.01,
            0.50,
            0.10,
            0.01,
            help="Sets the minimum relative peak height required for a signal to be labeled as a detected peak.",
        )
        min_distance = st.slider(
            "Minimum Peak Distance (m/z)",
            0.1,
            5.0,
            0.5,
            0.1,
            help="Prevents neighboring peaks from being counted separately unless they are at least this far apart in m/z.",
        )
        top_n = st.slider(
            "Top Matches",
            1,
            10,
            5,
            1,
            help="Controls how many reference-library matches are returned in the ranked comparison table.",
        )
        match_tolerance = st.slider(
            "Matching Tolerance (m/z)",
            0.2,
            3.0,
            1.0,
            0.1,
            help="Defines how close two peaks must be in m/z to count as a match during library scoring.",
        )

    file_source = None
    file_label = None
    export_row = {
        "microbe": "Uploaded Sample",
        "organism_type": "",
        "gram_status": "",
        "analyte_category": "",
        "strain": "",
        "condition": "",
        "replicate_id": "",
    }

    if uploaded is not None:
        file_source = BytesIO(uploaded.getvalue())
        file_label = uploaded.name
        export_row["microbe"] = uploaded.name
    elif selected_path is not None and selected_path.exists():
        file_source = str(selected_path)
        file_label = selected_path.name
        export_row["microbe"] = selected_path.name

    if file_source is None:
        st.info("Upload an mzXML file or choose one from the library.")
        return

    parsed = spectrum_from_mzxml(
        file_source,
        rel_height=rel_height,
        min_distance_mz=min_distance,
    )

    if parsed is None:
        st.error("No usable scans were found in this mzXML file.")
        return

    st.markdown("### Spectrum Preview")
    st.write(f"Representative Scan: {parsed['scan_num']}")

    fig = plot_mzxml_spectrum(
        parsed["mzs"],
        parsed["intensities"],
        parsed["peak_mz"],
        parsed["peak_i"],
        title=file_label,
    )
    st.pyplot(fig)

    with st.expander("Show Detected Peaks"):
        peak_df = pd.DataFrame({
            "mz": parsed["peak_mz"],
            "intensity": parsed["peak_i"],
        }).sort_values("intensity", ascending=False)
        peak_df = prettify_dataframe_columns(peak_df)
        st.dataframe(peak_df, use_container_width=True)

    ppt_path = create_ppt_export(export_row, parsed)
    with open(ppt_path, "rb") as f:
        st.download_button(
            label="Export Uploaded Spectrum to PowerPoint",
            data=f,
            file_name=f"{normalize_search_text(file_label).replace(' ', '_')}_spectrum.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="upload_ppt_export",
            help="Download a PowerPoint summary slide for the uploaded or selected spectrum.",
        )

    st.markdown("---")
    st.markdown("### Closest Library Matches")

    predicted_class = infer_profile_class(parsed)
    st.write(f"Predicted Profile Class: {predicted_class}")

    refs = load_reference_library(metadata)
    match_df = compare_to_library(parsed, refs, top_n=top_n, tolerance=match_tolerance)

    if match_df.empty:
        st.warning("No reference mzXML files were found in the metadata library.")
    else:
        st.dataframe(match_df, use_container_width=True)


def main():
    st.title("Microbial Reference Library")
    st.caption("Interactive microbial reference spectra browser and mzXML matching workflow")

    metadata = load_metadata()

    if metadata.empty:
        st.error("data/metadata.csv was not found or is empty.")
        st.stop()

    tab1, tab2 = st.tabs(["Library Browser", "mzXML Upload and Matching"])

    with tab1:
        render_library_browser(metadata)

    with tab2:
        render_mzxml_matcher(metadata)

    st.markdown("---")
    st.markdown(
        """
        <div style='text-align: center; padding: 10px; font-size: 13px; color: #9aa0a6;'>
            Developed by <b>modata-analytics</b> |
            <a href="https://github.com/modata-analytics" target="_blank" style="color: #58a6ff;">
                View GitHub
            </a>
        </div>
        """,
        unsafe_allow_html=True,
    )


if __name__ == "__main__":
    main()